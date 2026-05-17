#!/usr/bin/env python3
"""
模板与数据文件兼容性诊断工具
检查 Excel 列名和模板占位符是否匹配
"""
import sys
import re
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent / "src"))

from pptx import Presentation
from docx import Document
from excel_reader import ExcelReader


def read_excel_headers(filepath):
    reader = ExcelReader(filepath)
    try:
        reader.read()
        return reader.headers
    except Exception as e:
        print(f"    [WARN] 读取 Excel 失败: {e}")
        return []


def find_pptx_placeholders(filepath):
    """Find all {{...}} and {...} placeholders in PPTX text"""
    prs = Presentation(filepath)
    placeholders = set()
    placeholder_pattern = re.compile(r"\{\{[^}]+\}\}|\{[^}]+\}")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text
                for m in placeholder_pattern.finditer(text):
                    placeholders.add(m.group())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text
                        for m in placeholder_pattern.finditer(text):
                            placeholders.add(m.group())

    return sorted(placeholders)


def find_docx_placeholders(filepath):
    """Find all {{...}} and {...} placeholders in DOCX"""
    doc = Document(filepath)
    placeholders = set()
    placeholder_pattern = re.compile(r"\{\{[^}]+\}\}|\{[^}]+\}")

    all_text = []
    for para in doc.paragraphs:
        all_text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    all_text.append(para.text)

    for text in all_text:
        for m in placeholder_pattern.finditer(text):
            placeholders.add(m.group())

    return sorted(placeholders)


def extract_key(placeholder):
    """Extract the key from a placeholder like {{name}} -> name"""
    m = re.match(r"\{\{([^}]+)\}\}", placeholder)
    if m:
        return m.group(1).strip()
    m = re.match(r"\{([^}]+)\}", placeholder)
    if m:
        return m.group(1).strip()
    return placeholder.strip()


def main():
    import argparse
    parser = argparse.ArgumentParser(description="诊断模板与数据文件兼容性")
    parser.add_argument("excel", help="Excel 数据文件路径")
    parser.add_argument("--pptx", help="PPTX 模板路径")
    parser.add_argument("--docx", help="DOCX 模板路径")
    args = parser.parse_args()

    print("=" * 60)
    print("  模板与数据文件兼容性诊断")
    print("=" * 60)

    # 1. Read Excel headers
    print(f"\n[1] Excel 列名 ({args.excel}):")
    headers = read_excel_headers(args.excel)
    if headers:
        for i, h in enumerate(headers, 1):
            print(f"    {i}. [{h}]")
    else:
        print("    [WARN] 未找到表头（第二行）")
        return

    # Also show the hardcoded mappings
    print(f"\n[2] 硬编码字段映射:")
    mappings = {
        "姓名": ["{{姓名}}", "{姓名}", "{{name}}", "{name}"],
        "单位": ["{{单位}}", "{单位}", "{{unit}}", "{unit}"],
        "职务": ["{{职务}}", "{职务}", "{{position}}", "{position}"],
    }
    for field, patterns in mappings.items():
        exists = "[OK]" if field in headers else "[MISSING]"
        print(f"    {exists} {field}: {', '.join(patterns)}")

    # 3. Check PPTX placeholders
    if args.pptx:
        print(f"\n[3] PPTX 模板占位符 ({args.pptx}):")
        phs = find_pptx_placeholders(args.pptx)
        if phs:
            for ph in phs:
                key = extract_key(ph)
                match = "[OK]" if key in headers else "[MISSING]"
                print(f"    {match} {ph} -> [{key}]")
        else:
            print("    [WARN] 未找到任何 {{...}} 或 {...} 占位符！")
            print("    请确认模板中的占位符格式正确。")

    # 4. Check DOCX placeholders
    if args.docx:
        print(f"\n[4] DOCX 模板占位符 ({args.docx}):")
        phs = find_docx_placeholders(args.docx)
        if phs:
            for ph in phs:
                key = extract_key(ph)
                match = "[OK]" if key in headers else "[MISSING]"
                print(f"    {match} {ph} -> [{key}]")
        else:
            print("    [WARN] 未找到任何 {{...}} 或 {...} 占位符！")
            print("    请确认模板中的占位符格式正确。")

    # 5. Summary
    print(f"\n[5] 诊断结论:")
    all_phs = set()
    if args.pptx:
        all_phs.update(extract_key(p) for p in find_pptx_placeholders(args.pptx))
    if args.docx:
        all_phs.update(extract_key(p) for p in find_docx_placeholders(args.docx))

    if not all_phs:
        print("    [ERROR] 模板中没有找到占位符！")
        print("    占位符格式必须是 {{列名}} 或 {列名}")
        print("    例如: {{姓名}}、{{单位}}、{name}")
        return

    unmatched = [k for k in all_phs if k not in headers]
    matched = [k for k in all_phs if k in headers]

    if matched:
        print(f"    [OK] 匹配的占位符: {matched}")
    if unmatched:
        print(f"    [MISMATCH] 无匹配的占位符: {unmatched}")
        print(f"    这些占位符在 Excel 中找不到对应列名。")
        print(f"    请在 Excel 中添加对应列，或修改模板占位符为已有列名。")

    if not unmatched and matched:
        print("    [OK] 所有占位符都有对应的 Excel 列！模板应该能正常替换。")

    print("\n" + "=" * 60)


if __name__ == "__main__":
    main()
