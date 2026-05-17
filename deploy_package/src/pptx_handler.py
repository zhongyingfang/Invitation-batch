"""
PPTX处理模块
批量填充PPTX邀请函模板并生成高清PNG图片
"""
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from typing import Dict, Any
from PIL import Image
from utils import find_soffice
import io


class PPTXHandler:
    """处理PPTX文档"""
    
    def __init__(self, template_path: str, output_dir: str):
        """初始化PPTX处理器
        
        Args:
            template_path: PPTX模板文件路径
            output_dir: 输出目录路径
        """
        self.template_path = template_path
        self.output_dir = output_dir
        self.placeholder_keywords = ["姓名", "名字", "{{name}}", "{name}", "NAME", "{{"]
    
    def fill_template(self, data: Dict[str, Any], output_name: str) -> str:
        """使用数据填充PPTX模板
        
        Args:
            data: 包含填充数据的字典
            output_name: 输出文件名（不含扩展名）
        
        Returns:
            生成的PPTX文件路径
        """
        try:
            # 加载模板
            prs = Presentation(self.template_path)
            
            def replace_placeholders(text: str) -> str:
                if text is None:
                    return text
                replaced = str(text)
                # Chinese key placeholders
                replaced = replaced.replace("{{姓名}}", str(data.get("姓名", "")))
                replaced = replaced.replace("{姓名}", str(data.get("姓名", "")))
                replaced = replaced.replace("{{单位}}", str(data.get("单位", "")))
                replaced = replaced.replace("{单位}", str(data.get("单位", "")))
                replaced = replaced.replace("{{职务}}", str(data.get("职务", "")))
                replaced = replaced.replace("{职务}", str(data.get("职务", "")))
                # Lowercase hardcoded placeholders
                replaced = replaced.replace("{{name}}", str(data.get("姓名", "")))
                replaced = replaced.replace("{name}", str(data.get("姓名", "")))
                replaced = replaced.replace("{{unit}}", str(data.get("单位", "")))
                replaced = replaced.replace("{unit}", str(data.get("单位", "")))
                replaced = replaced.replace("{{position}}", str(data.get("职务", "")))
                replaced = replaced.replace("{position}", str(data.get("职务", "")))

                for key, value in data.items():
                    if key is None:
                        continue
                    if key not in ("姓名", "单位", "职务", "name", "unit", "position"):
                        replaced = replaced.replace(f"{{{{{key}}}}}", str(value or ""))
                        replaced = replaced.replace(f"{{{key}}}", str(value or ""))
                return replaced

            def replace_placeholders_in_runs(paragraph):
                """跨run替换占位符，保留样式和空格"""
                runs = list(paragraph.runs)
                if not runs:
                    return

                full_text = ''.join(run.text for run in runs)
                replaced_full = replace_placeholders(full_text)
                if replaced_full == full_text:
                    return

                # 找出所有占位符及其在full_text中的位置
                placeholders_to_replace = []
                # Chinese key placeholders
                chinese_patterns = {
                    "姓名": str(data.get("姓名", "")),
                    "单位": str(data.get("单位", "")),
                    "职务": str(data.get("职务", "")),
                }
                for ph_key, ph_value in chinese_patterns.items():
                    for pattern in [f"{{{{{ph_key}}}}}", f"{{{ph_key}}}"]:
                        if pattern in full_text and pattern != ph_value:
                            placeholders_to_replace.append((pattern, ph_value))

                # English hardcoded placeholders
                hardcoded_patterns = {
                    "name": str(data.get("姓名", "")),
                    "unit": str(data.get("单位", "")),
                    "position": str(data.get("职务", "")),
                }
                for ph_key, ph_value in hardcoded_patterns.items():
                    for pattern in [f"{{{{{ph_key}}}}}", f"{{{ph_key}}}"]:
                        if pattern in full_text and pattern != ph_value:
                            placeholders_to_replace.append((pattern, ph_value))

                # 数据键的占位符
                for key in data.keys():
                    if key in ("姓名", "单位", "职务", "name", "unit", "position"):
                        continue
                    for pattern in [f"{{{{{key}}}}}", f"{{{key}}}"]:
                        if pattern in full_text and pattern != str(data.get(key, "")):
                            placeholders_to_replace.append((pattern, str(data.get(key, "") or "")))

                if not placeholders_to_replace:
                    return

                # 收集所有占位符的(start, end, replacement)
                raw_replacements = []
                for ph, repl in placeholders_to_replace:
                    start = full_text.find(ph)
                    end = start + len(ph)
                    raw_replacements.append((start, end, repl))

                # 按起始位置排序，去除重叠的占位符（保留较长的，处理嵌套情况）
                raw_replacements.sort(key=lambda x: x[0])
                filtered = []
                for start, end, repl in raw_replacements:
                    # 检查是否与已有占位符重叠
                    contained = False
                    for f_start, f_end, f_repl in filtered:
                        if start >= f_start and end <= f_end:
                            # 当前占位符被已有占位符完全包含，跳过
                            contained = True
                            break
                        elif start < f_end and end > f_start:
                            # 有交叉重叠，保留较长的
                            if (end - start) > (f_end - f_start):
                                filtered.remove((f_start, f_end, f_repl))
                                filtered.append((start, end, repl))
                            contained = True
                            break
                    if not contained:
                        filtered.append((start, end, repl))

                # 保存原始run文本（永远不修改）
                original_texts = [run.text for run in runs]

                # 从右往左替换，只使用原始位置，不使用offset
                replacements = sorted(filtered, key=lambda x: x[0], reverse=True)

                for ph_start, ph_end, replacement in replacements:
                    # 使用原始位置找到占位符所在的run范围
                    pos = 0
                    start_run_idx = None
                    end_run_idx = None
                    for i, orig_text in enumerate(original_texts):
                        run_end = pos + len(orig_text)
                        if start_run_idx is None and pos <= ph_start < run_end:
                            start_run_idx = i
                        if end_run_idx is None and pos < ph_end <= run_end:
                            end_run_idx = i
                        pos = run_end

                    if start_run_idx is None:
                        start_run_idx = 0
                    if end_run_idx is None:
                        end_run_idx = start_run_idx

                    # 计算在start_run中的局部位置（使用原始文本）
                    run_start_pos = sum(len(t) for t in original_texts[:start_run_idx])
                    local_ph_start = ph_start - run_start_pos
                    local_ph_end = ph_end - run_start_pos

                    # 修改start_run：替换占位符部分
                    first_run = runs[start_run_idx]
                    before = original_texts[start_run_idx][:local_ph_start] if local_ph_start > 0 else ""
                    if start_run_idx == end_run_idx:
                        # 占位符在同一个run内
                        after = original_texts[start_run_idx][local_ph_end:] if local_ph_end < len(original_texts[start_run_idx]) else ""
                        first_run.text = before + replacement + after
                    else:
                        # 占位符跨多个run
                        run_end_pos = sum(len(t) for t in original_texts[:end_run_idx])
                        local_ph_end_in_last = ph_end - run_end_pos
                        after = original_texts[end_run_idx][local_ph_end_in_last:] if local_ph_end_in_last < len(original_texts[end_run_idx]) else ""

                        first_run.text = before + replacement + after
                        # 清空中间的所有run
                        for j in range(start_run_idx + 1, end_run_idx + 1):
                            runs[j].text = ""

            def adjust_name_line(slide, name_value: str):
                if not name_value:
                    return

                line_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE]
                if not line_shapes:
                    return

                text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape.text_frame is not None and name_value in shape.text_frame.text]
                if not text_shapes:
                    return

                text_shape = text_shapes[0]
                line = line_shapes[0]

                font_size_emu = None
                name_font_color = None
                for paragraph in text_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if name_value in run.text:
                            if run.font.size is not None:
                                font_size_emu = run.font.size
                            try:
                                if run.font.color and run.font.color.rgb:
                                    name_font_color = run.font.color.rgb
                            except Exception:
                                pass
                            break
                    if font_size_emu:
                        break

                if font_size_emu is None:
                    font_size_emu = Pt(12)

                name_width_emu = sum(
                    int(font_size_emu * (1.0 if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f' or '\uff00' <= c <= '\uffef' else 0.5))
                    for c in name_value
                )

                text = text_shape.text or ''
                name_pos = text.find(name_value)

                if name_pos >= 0:
                    prefix = text[:name_pos]
                    prefix_width_emu = sum(
                        int(font_size_emu * (1.0 if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f' or '\uff00' <= c <= '\uffef' else 0.5))
                        for c in prefix
                    )
                    shift_emu = font_size_emu // 2
                    line.left = int(text_shape.left + prefix_width_emu - shift_emu)

                line.width = int(name_width_emu + font_size_emu)
                if name_font_color:
                    try:
                        line.line.color.rgb = name_font_color
                    except Exception:
                        pass

            # 遍历所有幻灯片
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        for paragraph in shape.text_frame.paragraphs:
                            replace_placeholders_in_runs(paragraph)

                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    replace_placeholders_in_runs(paragraph)

                adjust_name_line(slide, str(data.get("姓名", "")).strip())
            
            # 保存PPTX文件
            output_path = os.path.join(self.output_dir, f"{output_name}.pptx")
            os.makedirs(self.output_dir, exist_ok=True)
            prs.save(output_path)
            
            print(f"生成PPTX文件: {output_path}")
            return output_path
            
        except Exception as e:
            print(f"填充PPTX模板失败: {e}")
            raise
    
    def convert_to_png(self, pptx_path: str, png_output_dir: str, dpi: int = 300) -> list:
        """将PPTX转换为高清PNG图片
        
        Args:
            pptx_path: PPTX文件路径
            png_output_dir: PNG输出目录
            dpi: 分辨率（dpi）
        
        Returns:
            生成的PNG文件路径列表
        """
        try:
            import subprocess
            
            os.makedirs(png_output_dir, exist_ok=True)
            
            # 获取文件名（不含扩展名）
            base_name = Path(pptx_path).stem
            
            # 使用LibreOffice将PPTX转换为PDF（中间步骤）
            pdf_path = os.path.join(png_output_dir, f"{base_name}.pdf")
            
            soffice_exe = find_soffice()
            if not soffice_exe:
                print("未找到 LibreOffice，跳过 PNG 转换。（安装 LibreOffice 或设置 LO_PATH 环境变量）")
                return []

            try:
                cmd = [
                    soffice_exe,
                    "--headless",
                    "--norestore",
                    "--invisible",
                    "--convert-to", "pdf",
                    "--outdir", png_output_dir,
                    pptx_path
                ]
                subprocess.run(cmd, check=True, capture_output=True, timeout=120)
                print(f"已转换为PDF: {pdf_path}")
            except Exception as e:
                print(f"LibreOffice转换失败: {e}，尝试使用pdf2image的替代方案")
            
            # 使用pdf2image将PDF转换为PNG
            if os.path.exists(pdf_path):
                import glob
                from pdf2image import convert_from_path
                from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError

                # 在 Linux/Docker 环境中显式设置 pdftoppm 路径
                poppler_path = None
                if sys.platform != 'win32':
                    for candidate in ['/usr/bin/pdftoppm', '/usr/local/bin/pdftoppm']:
                        if os.path.isfile(candidate):
                            poppler_path = os.path.dirname(candidate)
                            break
                    if not poppler_path:
                        import subprocess as _sp
                        try:
                            result = _sp.run(['which', 'pdftoppm'], capture_output=True, text=True)
                            if result.returncode == 0 and result.stdout.strip():
                                poppler_path = os.path.dirname(result.stdout.strip())
                        except Exception:
                            pass

                try:
                    images = convert_from_path(pdf_path, dpi=dpi, poppler_path=poppler_path)
                except PDFInfoNotInstalledError:
                    print("poppler-utils 未安装，跳过 PNG 转换")
                    return []
                
                png_paths = []
                
                for i, image in enumerate(images, 1):
                    png_path = os.path.join(png_output_dir, f"{base_name}_page{i}.png")
                    image.save(png_path, 'PNG')
                    png_paths.append(png_path)
                    print(f"生成PNG文件: {png_path}")
                
                return png_paths
            else:
                print(f"PDF文件不存在，无法转换为PNG")
                return []
            
        except Exception as e:
            print(f"转换为PNG失败: {e}")
            # 返回空列表而不抛出异常，允许程序继续
            return []


if __name__ == "__main__":
    handler = PPTXHandler("invitation_template.pptx", "output_documents")
    data = {
        "姓名": "张三",
        "单位": "测试公司",
        "职务": "总经理"
    }
    pptx_file = handler.fill_template(data, "邀请函_张三")
    # handler.convert_to_png(pptx_file, "output_images")
